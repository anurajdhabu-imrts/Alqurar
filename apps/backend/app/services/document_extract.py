"""Extract plain text from uploaded construction documents.

Supports PDF (PyMuPDF, falling back to pypdf), Word (.docx, python-docx),
PowerPoint (.pptx, python-pptx), Excel (.xlsx, openpyxl) and plain-text formats.
Image files and image-only ("scanned") PDFs yield no embedded text here — those
are handled by OCR (see app.services.ai_service.ocr_document). Anything else falls
back to the filename only, so the AI can still classify from the name.
"""

import io

# Keep extracted text bounded so a huge file can't blow up the token bill.
MAX_CHARS = 60_000
# Cap pages parsed from a PDF — a large or scanned PDF can otherwise take a long
# time for little useful text. MAX_CHARS is reached long before this for text PDFs.
MAX_PDF_PAGES = 200

_TEXT_EXTS = {"txt", "md", "csv", "log", "xml", "json", "rtf"}
_IMAGE_EXTS = {"png", "jpg", "jpeg", "tif", "tiff", "webp", "gif", "bmp"}


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def is_ocr_candidate(filename: str) -> bool:
    """True for formats whose text may need OCR (PDF or image)."""
    ext = _ext(filename)
    return ext == "pdf" or ext in _IMAGE_EXTS


def _extract_pdf(raw: bytes, max_chars: int = MAX_CHARS) -> str:
    """Fast PDF text extraction via PyMuPDF, with a pypdf fallback."""
    try:
        import fitz  # PyMuPDF

        parts, total = [], 0
        with fitz.open(stream=raw, filetype="pdf") as doc:
            for page in doc[:MAX_PDF_PAGES] if len(doc) > MAX_PDF_PAGES else doc:
                page_text = page.get_text("text") or ""
                parts.append(page_text)
                total += len(page_text)
                if total >= max_chars:
                    break
        return "\n".join(parts)
    except Exception:
        # Fallback to pypdf if PyMuPDF isn't available or chokes on the file.
        from pypdf import PdfReader

        reader = PdfReader(io.BytesIO(raw))
        parts, total = [], 0
        for page in reader.pages[:MAX_PDF_PAGES]:
            page_text = page.extract_text() or ""
            parts.append(page_text)
            total += len(page_text)
            if total >= max_chars:
                break
        return "\n".join(parts)


def _extract_docx(raw: bytes) -> str:
    """Extract text from a .docx, INCLUDING tables, headers and footers.

    Particular Conditions / Contract Data are frequently laid out in tables (and
    sometimes headers), which `document.paragraphs` alone misses entirely — so we
    walk the body in document order and pull every table cell as well, then append
    header/footer text from each section.
    """
    import docx
    from docx.document import Document as _Doc
    from docx.oxml.table import CT_Tbl
    from docx.oxml.text.paragraph import CT_P
    from docx.table import Table, _Cell
    from docx.text.paragraph import Paragraph

    document = docx.Document(io.BytesIO(raw))

    def _iter_block_text(parent) -> list[str]:
        """Yield text from paragraphs and tables of a body/cell, in order."""
        if isinstance(parent, _Doc):
            elm = parent.element.body
        elif isinstance(parent, _Cell):
            elm = parent._tc
        else:
            elm = parent
        out: list[str] = []
        for child in elm.iterchildren():
            if isinstance(child, CT_P):
                t = Paragraph(child, parent).text
                if t.strip():
                    out.append(t)
            elif isinstance(child, CT_Tbl):
                table = Table(child, parent)
                for row in table.rows:
                    cells = []
                    for cell in row.cells:
                        cells.append(" ".join(_iter_block_text(cell)).strip())
                    line = " | ".join(c for c in cells if c)
                    if line.strip():
                        out.append(line)
        return out

    parts = _iter_block_text(document)

    # Headers and footers (subject-line boilerplate, appendix titles, etc.).
    for section in document.sections:
        for hf in (section.header, section.footer):
            try:
                hf_text = "\n".join(p.text for p in hf.paragraphs if p.text.strip())
            except Exception:
                hf_text = ""
            if hf_text.strip():
                parts.append(hf_text)

    return "\n".join(parts)


def extract_text(filename: str, raw: bytes, max_chars: int = MAX_CHARS) -> tuple[str, bool]:
    """Return (text, truncated). `text` is "" when the format isn't extractable
    (e.g. an image, or a scanned PDF with no embedded text → handled by OCR).

    `max_chars` bounds the extracted text. It defaults to MAX_CHARS (fine for the
    short per-document analysis), but callers reading a long document end-to-end —
    e.g. extracting clauses from a full contract — pass a larger value so the
    important later sections aren't cut off."""
    ext = _ext(filename)
    text = ""

    try:
        if ext == "pdf":
            text = _extract_pdf(raw, max_chars)
        elif ext in ("docx", "doc"):
            # Handles paragraphs AND tables/headers/footers. A legacy binary .doc
            # is not a zip and raises here → caught below → OCR/empty fallback.
            text = _extract_docx(raw)
        elif ext == "pptx":
            from pptx import Presentation

            prs = Presentation(io.BytesIO(raw))
            chunks = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        chunks.append(shape.text_frame.text)
                    if shape.has_table:
                        for row in shape.table.rows:
                            chunks.append(" | ".join(c.text for c in row.cells))
            text = "\n".join(c for c in chunks if c)
        elif ext in ("xlsx", "xlsm"):
            from openpyxl import load_workbook

            wb = load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            rows = []
            for ws in wb.worksheets:
                rows.append(f"# Sheet: {ws.title}")
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) for c in row if c is not None]
                    if cells:
                        rows.append(" | ".join(cells))
            text = "\n".join(rows)
        elif ext in _TEXT_EXTS:
            text = raw.decode("utf-8", errors="replace")
    except Exception:
        # Corrupt/unsupported encoding — fall back to filename-only analysis.
        text = ""

    text = text.strip()
    truncated = len(text) > max_chars
    if truncated:
        text = text[:max_chars]
    return text, truncated
